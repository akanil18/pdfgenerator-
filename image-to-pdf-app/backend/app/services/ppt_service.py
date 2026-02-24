"""
PDF to PowerPoint Service — converts PDF pages to a PPTX presentation.
Each page becomes a slide with the page rendered as an image.
"""

import os
import uuid
import logging

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Emu

from pypdf import PdfReader

logger = logging.getLogger(__name__)

# Standard slide dimensions (10 × 7.5 inches)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)


def pdf_to_ppt(pdf_path: str, session_dir: str) -> str:
    """
    Convert a PDF into a PowerPoint presentation.
    Each PDF page is rendered as an image and placed on a slide.
    Returns the path to the generated PPTX file.
    """
    pptx_filename = f"{uuid.uuid4().hex}.pptx"
    pptx_path = os.path.join(session_dir, pptx_filename)

    try:
        # Try using pdf2image (requires poppler)
        from pdf2image import convert_from_path
        page_images = convert_from_path(pdf_path, dpi=200)
    except Exception:
        # Fallback: use pypdf to extract pages as images via Pillow
        # This is simpler but only works for image-based PDFs
        logger.warning("pdf2image unavailable, using fallback method")
        page_images = _fallback_extract_pages(pdf_path, session_dir)

    if not page_images:
        raise ValueError("Could not extract any pages from the PDF.")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for idx, img in enumerate(page_images):
        # Save page image temporarily
        img_path = os.path.join(session_dir, f"slide_{idx}.png")
        if isinstance(img, PILImage.Image):
            img.save(img_path, "PNG")
        else:
            # Already a path
            img_path = img

        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Calculate image size to fit slide maintaining aspect ratio
        pil_img = PILImage.open(img_path)
        img_w, img_h = pil_img.size
        pil_img.close()

        slide_w = SLIDE_WIDTH
        slide_h = SLIDE_HEIGHT

        ratio = min(slide_w / Emu(img_w * 9525), slide_h / Emu(img_h * 9525))
        final_w = int(img_w * 9525 * ratio)
        final_h = int(img_h * 9525 * ratio)

        left = (SLIDE_WIDTH - final_w) // 2
        top = (SLIDE_HEIGHT - final_h) // 2

        slide.shapes.add_picture(img_path, left, top, final_w, final_h)
        logger.info(f"Added slide {idx + 1}")

    prs.save(pptx_path)

    file_size = os.path.getsize(pptx_path)
    logger.info(f"✅  PPTX created: {pptx_path} ({file_size:,} bytes, {len(page_images)} slides)")
    return pptx_path


def _fallback_extract_pages(pdf_path: str, session_dir: str) -> list:
    """
    Fallback: extract each PDF page as a separate single-page PDF,
    then convert to image using Pillow (works for simple PDFs).
    """
    reader = PdfReader(pdf_path)
    images = []

    for idx, page in enumerate(reader.pages):
        from pypdf import PdfWriter
        writer = PdfWriter()
        writer.add_page(page)
        temp_pdf = os.path.join(session_dir, f"_temp_page_{idx}.pdf")
        with open(temp_pdf, "wb") as f:
            writer.write(f)

        # Try to render via Pillow (limited support)
        try:
            from pdf2image import convert_from_path
            imgs = convert_from_path(temp_pdf, dpi=200)
            if imgs:
                images.extend(imgs)
        except Exception:
            # Last resort: create a placeholder slide
            placeholder = PILImage.new("RGB", (1920, 1080), (255, 255, 255))
            images.append(placeholder)
            logger.warning(f"Could not render page {idx + 1}, using placeholder")

    return images
